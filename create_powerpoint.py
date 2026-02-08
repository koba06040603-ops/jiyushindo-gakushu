#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
全46枚のスライドをPowerPointにまとめるスクリプト
"""

from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO
from PIL import Image
import time

# スライドデータ（透かしなし版）
slides_data = [
    # 第1部（1-15）
    {'num': 1, 'title': 'タイトルスライド', 'url': 'https://www.genspark.ai/api/files/s/UjpjgWGy?cache_control=3600'},
    {'num': 2, 'title': '現状の課題', 'url': 'https://www.genspark.ai/api/files/s/1QVsdgI2?cache_control=3600'},
    {'num': 3, 'title': '本システムの解決策', 'url': 'https://www.genspark.ai/api/files/s/nchqL6pn?cache_control=3600'},
    {'num': 4, 'title': '世界最先端の3つの強み', 'url': 'https://www.genspark.ai/api/files/s/zNoBn2E3?cache_control=3600'},
    {'num': 5, 'title': '12理論エビデンス一覧（F1-F4）', 'url': 'https://www.genspark.ai/api/files/s/p3ZVmyGZ?cache_control=3600'},
    {'num': 6, 'title': '12理論エビデンス一覧（F5-F8）', 'url': 'https://www.genspark.ai/api/files/s/OiT172mH?cache_control=3600'},
    {'num': 7, 'title': '12理論エビデンス一覧（F9-F12）', 'url': 'https://www.genspark.ai/api/files/s/OAvRW9bu?cache_control=3600'},
    {'num': '7A', 'title': '超高効果量研究トップ5', 'url': 'https://www.genspark.ai/api/files/s/dLE562yk?cache_control=3600'},
    {'num': 8, 'title': 'エビデンスの信頼性', 'url': 'https://www.genspark.ai/api/files/s/olzgwdg2?cache_control=3600'},
    {'num': 9, 'title': '導入効果の科学的予測', 'url': 'https://www.genspark.ai/api/files/s/i0rUJigi?cache_control=3600'},
    {'num': 10, 'title': '国際的評価と整合性', 'url': 'https://www.genspark.ai/api/files/s/zQd0f2Kw?cache_control=3600'},
    {'num': 11, 'title': 'システム構成（Gemini 3 Flash）', 'url': 'https://www.genspark.ai/api/files/s/q2UNaxq4?cache_control=3600'},
    {'num': 12, 'title': 'セキュリティとプライバシー', 'url': 'https://www.genspark.ai/api/files/s/k86CNWl0?cache_control=3600'},
    {'num': 13, 'title': '他システムとの比較', 'url': 'https://www.genspark.ai/api/files/s/91RbH6VG?cache_control=3600'},
    {'num': 14, 'title': '導入実績と将来展望', 'url': 'https://www.genspark.ai/api/files/s/FJkah4OV?cache_control=3600'},
    {'num': 15, 'title': '第1部まとめ', 'url': 'https://www.genspark.ai/api/files/s/KfpFIvuR?cache_control=3600'},
    
    # 第2部（16-30）
    {'num': 16, 'title': '第2部イントロダクション', 'url': 'https://www.genspark.ai/api/files/s/bO8G6yKV?cache_control=3600'},
    {'num': 17, 'title': 'シンプルで安全なログイン', 'url': 'https://www.genspark.ai/api/files/s/OQJfiI9K?cache_control=3600'},
    {'num': 18, 'title': '学習開始画面', 'url': 'https://www.genspark.ai/api/files/s/2W172Ksm?cache_control=3600'},
    {'num': 19, 'title': '学習カード（3コース制）', 'url': 'https://www.genspark.ai/api/files/s/2W172Ksm?cache_control=3600'},
    {'num': 20, 'title': 'AI先生のソクラテス対話（Gemini 3 Flash）', 'url': 'https://www.genspark.ai/api/files/s/EDobbnfm?cache_control=3600'},
    {'num': 21, 'title': '学習スタイル自動判定', 'url': 'https://www.genspark.ai/api/files/s/nUAhT7W5?cache_control=3600'},
    {'num': 22, 'title': '児童用進捗ダッシュボード', 'url': 'https://www.genspark.ai/api/files/s/EDobbnfm?cache_control=3600'},
    {'num': 23, 'title': '教員用進捗ボード', 'url': 'https://www.genspark.ai/api/files/s/D2Wglfph?cache_control=3600'},
    {'num': 24, 'title': '間違いノートと復習', 'url': 'https://www.genspark.ai/api/files/s/avI78lbX?cache_control=3600'},
    {'num': 25, 'title': '学習計画表と振り返り', 'url': 'https://www.genspark.ai/api/files/s/EASfx4gY?cache_control=3600'},
    {'num': 26, 'title': '保護者向けレポート', 'url': 'https://www.genspark.ai/api/files/s/EgdygdwT?cache_control=3600'},
    {'num': 27, 'title': 'AI問題生成（Gemini 3 Flash）', 'url': 'https://www.genspark.ai/api/files/s/1pZxhP5u?cache_control=3600'},
    {'num': 28, 'title': '1秒以内に適応', 'url': 'https://www.genspark.ai/api/files/s/94FNIoNF?cache_control=3600'},
    {'num': 29, 'title': '予測分析とリスク検知', 'url': 'https://www.genspark.ai/api/files/s/Hf15A8X9?cache_control=3600'},
    {'num': 30, 'title': '第2部まとめ', 'url': 'https://www.genspark.ai/api/files/s/Nn4veBRi?cache_control=3600'},
    
    # 第3部（31-46）
    {'num': 31, 'title': '第3部イントロダクション', 'url': 'https://www.genspark.ai/api/files/s/J95Lws5o?cache_control=3600'},
    {'num': 32, 'title': '5ステップ導入プロセス', 'url': 'https://www.genspark.ai/api/files/s/HS2G4H9o?cache_control=3600'},
    {'num': 33, 'title': '事前準備の詳細', 'url': 'https://www.genspark.ai/api/files/s/2He0hfW5?cache_control=3600'},
    {'num': 34, 'title': '初期設定の詳細', 'url': 'https://www.genspark.ai/api/files/s/UmlBiwez?cache_control=3600'},
    {'num': 35, 'title': '教員研修の詳細', 'url': 'https://www.genspark.ai/api/files/s/n2Nt75bR?cache_control=3600'},
    {'num': 36, 'title': 'パイロット運用の詳細', 'url': 'https://www.genspark.ai/api/files/s/XPsbPcJ0?cache_control=3600'},
    {'num': 37, 'title': '本格導入の詳細', 'url': 'https://www.genspark.ai/api/files/s/zUGcqFM6?cache_control=3600'},
    {'num': 38, 'title': '成功のポイント1：教員の巻き込み', 'url': 'https://www.genspark.ai/api/files/s/by7iyMYa?cache_control=3600'},
    {'num': 39, 'title': '成功のポイント2：児童の動機づけ', 'url': 'https://www.genspark.ai/api/files/s/o7MP1idV?cache_control=3600'},
    {'num': 40, 'title': '成功のポイント3：保護者の理解', 'url': 'https://www.genspark.ai/api/files/s/tfqrYoVF?cache_control=3600'},
    {'num': 41, 'title': '成功のポイント4：定着サポート', 'url': 'https://www.genspark.ai/api/files/s/gXnub3Nr?cache_control=3600'},
    {'num': 42, 'title': '成功のポイント5：PDCAサイクル', 'url': 'https://www.genspark.ai/api/files/s/Clk4ZjGr?cache_control=3600'},
    {'num': 43, 'title': '効果測定と成果指標', 'url': 'https://www.genspark.ai/api/files/s/HFQbZPw6?cache_control=3600'},
    {'num': 44, 'title': '松川村での特別条件', 'url': 'https://www.genspark.ai/api/files/s/DYJqpHGK?cache_control=3600'},
    {'num': 45, 'title': 'まとめと次のアクション', 'url': 'https://www.genspark.ai/api/files/s/u3eVM2r0?cache_control=3600'},
    {'num': 46, 'title': 'よくある質問（Q&A）', 'url': 'https://www.genspark.ai/api/files/s/xT0Yl4GH?cache_control=3600'},
]

def download_image(url, max_retries=3):
    """画像をダウンロード"""
    for attempt in range(max_retries):
        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            return BytesIO(response.content)
        except Exception as e:
            print(f"  ⚠️  ダウンロード失敗（試行 {attempt + 1}/{max_retries}）: {e}")
            if attempt < max_retries - 1:
                time.sleep(2)
            else:
                return None
    return None

def create_powerpoint():
    """PowerPointファイルを作成"""
    print("📊 PowerPoint作成開始...")
    print(f"📝 総スライド数: {len(slides_data)}枚\n")
    
    # 16:9のプレゼンテーションを作成
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9の幅
    prs.slide_height = Inches(5.625)  # 16:9の高さ
    
    success_count = 0
    failed_slides = []
    
    for i, slide_info in enumerate(slides_data, 1):
        print(f"[{i:2d}/46] スライド{slide_info['num']}: {slide_info['title']}")
        
        # 画像をダウンロード
        image_stream = download_image(slide_info['url'])
        
        if image_stream is None:
            print(f"  ❌ スキップ（ダウンロード失敗）\n")
            failed_slides.append(slide_info['num'])
            continue
        
        try:
            # 空白のスライドを追加
            blank_slide_layout = prs.slide_layouts[6]  # 6 = 空白レイアウト
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 画像を追加（スライド全体に配置）
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height
            
            slide.shapes.add_picture(image_stream, left, top, width, height)
            
            success_count += 1
            print(f"  ✅ 追加完了\n")
            
        except Exception as e:
            print(f"  ❌ スライド追加失敗: {e}\n")
            failed_slides.append(slide_info['num'])
        
        # サーバー負荷軽減のため少し待機
        time.sleep(0.5)
    
    # 保存
    output_path = '/home/user/webapp/自由進度学習システム_全46枚.pptx'
    prs.save(output_path)
    
    print("\n" + "="*60)
    print("🎉 PowerPoint作成完了！")
    print("="*60)
    print(f"📁 保存先: {output_path}")
    print(f"✅ 成功: {success_count}枚")
    if failed_slides:
        print(f"❌ 失敗: {len(failed_slides)}枚 - スライド番号: {', '.join(map(str, failed_slides))}")
    else:
        print("❌ 失敗: なし")
    print(f"📊 合計: {success_count}/{len(slides_data)}枚")
    print("="*60)
    
    return output_path, success_count, failed_slides

if __name__ == "__main__":
    output_path, success, failed = create_powerpoint()
