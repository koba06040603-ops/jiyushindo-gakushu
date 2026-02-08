#!/usr/bin/env python3
"""
PowerPoint作成スクリプト（テキストのみ版）
全46枚のスライド情報をPowerPointファイルに出力します。
画像はダウンロードせず、各スライドにタイトルとダウンロードURLを記載します。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 全スライド情報（タイトルとURL）
slides_data = [
    # 第1部（1-15）
    {"title": "スライド1: タイトルスライド", "url": "https://www.genspark.ai/api/files/s/UjpjgWGy?cache_control=3600", "part": 1},
    {"title": "スライド2: 現状の課題", "url": "https://www.genspark.ai/api/files/s/1QVsdgI2?cache_control=3600", "part": 1},
    {"title": "スライド3: 本システムの解決策", "url": "https://www.genspark.ai/api/files/s/nchqL6pn?cache_control=3600", "part": 1},
    {"title": "スライド4: 世界最先端の3つの強み", "url": "https://www.genspark.ai/api/files/s/zNoBn2E3?cache_control=3600", "part": 1},
    {"title": "スライド5: 12理論エビデンス一覧（F1-F4）", "url": "https://www.genspark.ai/api/files/s/p3ZVmyGZ?cache_control=3600", "part": 1},
    {"title": "スライド6: 12理論エビデンス一覧（F5-F8）", "url": "https://www.genspark.ai/api/files/s/OiT172mH?cache_control=3600", "part": 1},
    {"title": "スライド7: 12理論エビデンス一覧（F9-F12）", "url": "https://www.genspark.ai/api/files/s/OAvRW9bu?cache_control=3600", "part": 1},
    {"title": "スライド7A: 超高効果量研究トップ5", "url": "https://www.genspark.ai/api/files/s/dLE562yk?cache_control=3600", "part": 1},
    {"title": "スライド8: エビデンスの信頼性", "url": "https://www.genspark.ai/api/files/s/olzgwdg2?cache_control=3600", "part": 1},
    {"title": "スライド9: 導入効果の科学的予測", "url": "https://www.genspark.ai/api/files/s/i0rUJigi?cache_control=3600", "part": 1},
    {"title": "スライド10: 国際的評価と整合性", "url": "https://www.genspark.ai/api/files/s/zQd0f2Kw?cache_control=3600", "part": 1},
    {"title": "スライド11: システム構成（Gemini 3 Flash）", "url": "https://www.genspark.ai/api/files/s/q2UNaxq4?cache_control=3600", "part": 1},
    {"title": "スライド12: セキュリティとプライバシー", "url": "https://www.genspark.ai/api/files/s/k86CNWl0?cache_control=3600", "part": 1},
    {"title": "スライド13: 他システムとの比較", "url": "https://www.genspark.ai/api/files/s/91RbH6VG?cache_control=3600", "part": 1},
    {"title": "スライド14: 導入実績と将来展望", "url": "https://www.genspark.ai/api/files/s/FJkah4OV?cache_control=3600", "part": 1},
    {"title": "スライド15: 第1部まとめ", "url": "https://www.genspark.ai/api/files/s/KfpFIvuR?cache_control=3600", "part": 1},
    
    # 第2部（16-30）
    {"title": "スライド16: 第2部イントロダクション", "url": "https://www.genspark.ai/api/files/s/vPMBrcYo?cache_control=3600", "part": 2},
    {"title": "スライド17: AI先生との個別対話", "url": "https://www.genspark.ai/api/files/s/Fvps5MdV?cache_control=3600", "part": 2},
    {"title": "スライド18: リアルタイム適応学習", "url": "https://www.genspark.ai/api/files/s/OZPQsxl0?cache_control=3600", "part": 2},
    {"title": "スライド19: 間隔反復学習", "url": "https://www.genspark.ai/api/files/s/H8u9hVGM?cache_control=3600", "part": 2},
    {"title": "スライド20: AI先生のソクラテス対話（Gemini 3 Flash）", "url": "https://www.genspark.ai/api/files/s/EDobbnfm?cache_control=3600", "part": 2},
    {"title": "スライド21: AI分析による詳細フィードバック", "url": "https://www.genspark.ai/api/files/s/bJm5LTGj?cache_control=3600", "part": 2},
    {"title": "スライド22: ゲーミフィケーション", "url": "https://www.genspark.ai/api/files/s/C0Rr9Odk?cache_control=3600", "part": 2},
    {"title": "スライド23: 間違いノート", "url": "https://www.genspark.ai/api/files/s/Pfs1U0Wx?cache_control=3600", "part": 2},
    {"title": "スライド24: 学習経路最適化", "url": "https://www.genspark.ai/api/files/s/5wJ15iXS?cache_control=3600", "part": 2},
    {"title": "スライド25: 教員支援機能", "url": "https://www.genspark.ai/api/files/s/gUcQGJjq?cache_control=3600", "part": 2},
    {"title": "スライド26: リアルタイム監視ダッシュボード", "url": "https://www.genspark.ai/api/files/s/TqnPT61m?cache_control=3600", "part": 2},
    {"title": "スライド27: AI問題生成（Gemini 3 Flash）", "url": "https://www.genspark.ai/api/files/s/1pZxhP5u?cache_control=3600", "part": 2},
    {"title": "スライド28: 包括的レポート", "url": "https://www.genspark.ai/api/files/s/w0NVcm6f?cache_control=3600", "part": 2},
    {"title": "スライド29: 実際のスクリーンショット", "url": "https://www.genspark.ai/api/files/s/0s1F37BZ?cache_control=3600", "part": 2},
    {"title": "スライド30: 第2部まとめ", "url": "https://www.genspark.ai/api/files/s/N0NKZcB0?cache_control=3600", "part": 2},
    
    # 第3部（31-46）
    {"title": "スライド31: 第3部イントロダクション", "url": "https://www.genspark.ai/api/files/s/NZYWbhHi?cache_control=3600", "part": 3},
    {"title": "スライド32: 5ステップ導入プロセス", "url": "https://www.genspark.ai/api/files/s/qQnE7Clc?cache_control=3600", "part": 3},
    {"title": "スライド33: 事前準備", "url": "https://www.genspark.ai/api/files/s/2He0hfW5?cache_control=3600", "part": 3},
    {"title": "スライド34: 初期設定", "url": "https://www.genspark.ai/api/files/s/UmlBiwez?cache_control=3600", "part": 3},
    {"title": "スライド35: 教員研修", "url": "https://www.genspark.ai/api/files/s/TDjih9d4?cache_control=3600", "part": 3},
    {"title": "スライド36: パイロット運用", "url": "https://www.genspark.ai/api/files/s/mcejUMSg?cache_control=3600", "part": 3},
    {"title": "スライド37: 本格導入", "url": "https://www.genspark.ai/api/files/s/zUGcqFM6?cache_control=3600", "part": 3},
    {"title": "スライド38: 成功のポイント1：教員の巻き込み", "url": "https://www.genspark.ai/api/files/s/by7iyMYa?cache_control=3600", "part": 3},
    {"title": "スライド39: 成功のポイント2：児童の動機づけ", "url": "https://www.genspark.ai/api/files/s/o7MP1idV?cache_control=3600", "part": 3},
    {"title": "スライド40: 成功のポイント3：保護者の理解", "url": "https://www.genspark.ai/api/files/s/RyEiHHh2?cache_control=3600", "part": 3},
    {"title": "スライド41: 成功のポイント4：定着サポート", "url": "https://www.genspark.ai/api/files/s/6Pn1NXQ5?cache_control=3600", "part": 3},
    {"title": "スライド42: 成功のポイント5：PDCA", "url": "https://www.genspark.ai/api/files/s/Clk4ZjGr?cache_control=3600", "part": 3},
    {"title": "スライド43: 効果測定と成果指標", "url": "https://www.genspark.ai/api/files/s/HFQbZPw6?cache_control=3600", "part": 3},
    {"title": "スライド44: 松川村の特別条件", "url": "https://www.genspark.ai/api/files/s/xTOa8YXl?cache_control=3600", "part": 3},
    {"title": "スライド45: まとめと次のアクション", "url": "https://www.genspark.ai/api/files/s/e6lUa5H5?cache_control=3600", "part": 3},
    {"title": "スライド46: よくある質問（Q&A）", "url": "https://www.genspark.ai/api/files/s/xT0Yl4GH?cache_control=3600", "part": 3},
]

def create_text_only_presentation():
    """テキストのみのPowerPointファイルを作成"""
    
    print("📊 PowerPoint作成開始（テキストのみ版）...")
    print(f"📝 総スライド数: {len(slides_data)}枚\n")
    
    # プレゼンテーションを作成
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # タイトルスライド（表紙）
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
    
    # タイトルテキストボックス
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "松川村AI個別最適化学習システム"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # サブタイトル
    subtitle_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "全46枚スライド（2026年2月版）"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 注意書き
    note_box = title_slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
    note_frame = note_box.text_frame
    note_frame.text = "※ 各スライドに画像URLが記載されています。URLから画像をダウンロードしてください。"
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(16)
    note_para.font.color.rgb = RGBColor(200, 0, 0)
    note_para.alignment = PP_ALIGN.CENTER
    
    print("✅ 表紙スライド作成完了\n")
    
    # 各部のセパレータと内容スライド
    current_part = 0
    slide_count = 0
    
    for idx, slide_data in enumerate(slides_data, 1):
        part = slide_data['part']
        
        # 新しい部が始まる場合、セパレータスライドを追加
        if part != current_part:
            current_part = part
            separator_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 部のタイトル
            part_title_box = separator_slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(11.333), Inches(1.5)
            )
            part_title_frame = part_title_box.text_frame
            part_title_frame.text = f"第{part}部"
            part_title_para = part_title_frame.paragraphs[0]
            part_title_para.font.size = Pt(60)
            part_title_para.font.bold = True
            part_title_para.font.color.rgb = RGBColor(0, 102, 204)
            part_title_para.alignment = PP_ALIGN.CENTER
            
            # 背景色
            background = separator_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(240, 248, 255)
            
            print(f"📂 第{part}部セパレータ作成完了")
        
        # 内容スライドを作成
        content_slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_count += 1
        
        # スライドタイトル
        title_box = content_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 51, 102)
        
        # プレースホルダーテキスト
        placeholder_box = content_slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.333), Inches(1.5)
        )
        placeholder_frame = placeholder_box.text_frame
        placeholder_frame.text = "【画像をここに挿入してください】"
        placeholder_para = placeholder_frame.paragraphs[0]
        placeholder_para.font.size = Pt(36)
        placeholder_para.font.color.rgb = RGBColor(150, 150, 150)
        placeholder_para.alignment = PP_ALIGN.CENTER
        
        # ダウンロードURL
        url_box = content_slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.333), Inches(2)
        )
        url_frame = url_box.text_frame
        url_frame.word_wrap = True
        
        # URL見出し
        p1 = url_frame.paragraphs[0]
        p1.text = "🔗 画像ダウンロードURL:"
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0, 102, 204)
        
        # URL本体
        p2 = url_frame.add_paragraph()
        p2.text = slide_data['url']
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0, 0, 255)
        p2.font.underline = True
        
        # 手順説明
        p3 = url_frame.add_paragraph()
        p3.text = "\n📌 手順: URLをブラウザで開く → 画像を右クリック → 「名前を付けて画像を保存」→ PowerPointに挿入"
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(100, 100, 100)
        
        # フッター（スライド番号）
        footer_box = content_slide.shapes.add_textbox(
            Inches(11.5), Inches(7), Inches(1.5), Inches(0.4)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = f"{slide_count}/{len(slides_data)}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = RGBColor(150, 150, 150)
        footer_para.alignment = PP_ALIGN.RIGHT
        
        print(f"[{slide_count:2d}/{len(slides_data)}] {slide_data['title']} ✅")
    
    # 保存
    output_path = "/home/user/webapp/松川村AI学習システム_全46枚スライド_URL版.pptx"
    prs.save(output_path)
    
    print(f"\n✅ PowerPoint作成完了！")
    print(f"📁 保存先: {output_path}")
    print(f"📊 総スライド数: {len(prs.slides)}枚（表紙 + 3部セパレータ + 内容{len(slides_data)}枚）")
    print(f"\n💡 使い方:")
    print(f"   1) このファイルをダウンロード")
    print(f"   2) 各スライドのURLから画像をダウンロード")
    print(f"   3) PowerPointで各スライドのプレースホルダーに画像を挿入")
    
    return output_path

if __name__ == "__main__":
    create_text_only_presentation()
